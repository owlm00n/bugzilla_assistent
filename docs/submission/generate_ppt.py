#!/usr/bin/env python3
"""
OWLMYCLAW 火山杯比赛作品演示PPT生成器
生成 17 页演示PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
SLIDE_WIDTH = Inches(13.333)  # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)

# Color Palette
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # Dark navy
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)    # Blue accent
ACCENT_CYAN = RGBColor(0x00, 0xB4, 0xD8)    # Cyan
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xB0)   # Green
ACCENT_ORANGE = RGBColor(0xF4, 0x96, 0x3E)  # Orange/warning
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)     # Red
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DIM_GRAY = RGBColor(0x88, 0x88, 0x88)
CARD_BG = RGBColor(0x25, 0x25, 0x3A)

OUTPUT_PATH = "E:/go/04_AI/bugzilla_assistent/docs/submission/演示PPT.pptx"

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def set_slide_bg(slide, color):
    """Set solid background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=16, color=WHITE, bold=False,
                          alignment=PP_ALIGN.LEFT, line_spacing=1.5,
                          font_name="Microsoft YaHei"):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return tf


def add_card(slide, left, top, width, height, color=CARD_BG):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_bar(slide, title_text, subtitle_text=""):
    """Standard slide title bar."""
    # Top accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Title
    add_textbox(slide, 0.8, 0.3, 11, 0.7, title_text,
                font_size=36, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, 0.8, 1.0, 11, 0.5, subtitle_text,
                    font_size=16, color=LIGHT_GRAY)

    # Divider line
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.55),
        Inches(2), Inches(0.04)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = ACCENT_CYAN
    div.line.fill.background()


def add_page_number(slide, num, total=17):
    """Add page number at bottom right."""
    add_textbox(slide, 11.5, 7.0, 1.5, 0.4,
                f"{num}/{total}", font_size=10, color=DIM_GRAY,
                alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_BG)

# Decorative top bar
bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    SLIDE_WIDTH, Inches(0.15)
)
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# Main title
add_textbox(slide, 1.0, 1.5, 11, 1.2,
            "OWLMYCLAW", font_size=72, bold=True, color=WHITE)

# Subtitle
add_textbox(slide, 1.0, 2.8, 11, 0.8,
            "底层驱动全栈自动化研发助手 — AI 驱动的 USB/PD 调试工作流",
            font_size=22, color=ACCENT_CYAN)

# Divider
div = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(1.0), Inches(3.8),
    Inches(3), Inches(0.05)
)
div.fill.solid()
div.fill.fore_color.rgb = ACCENT_CYAN
div.line.fill.background()

# Info
add_multiline_textbox(slide, 1.0, 4.2, 11, 1.5, [
    "2026 火山杯 · 紫光展锐 AI 创新大赛",
    "研发与创新赛道",
    "提交日期: 2026-06-07"
], font_size=18, color=LIGHT_GRAY, line_spacing=1.8)


# ============================================================
# SLIDE 2: TABLE OF CONTENTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "目录", "CONTENTS")
add_page_number(slide, 2)

toc_items = [
    ("01", "课题背景", "USB/PD 驱动开发面临的挑战"),
    ("02", "痛点分析", "4大核心痛点深度剖析"),
    ("03", "解决方案总览", "OWLMYCLAW 6大模块架构"),
    ("04", "模块1: Bug智能助手", "Bug查询/解析/搜索/MCP服务"),
    ("05", "模块2: 日志分析器", "角色感知分析/增强查看器"),
    ("06", "模块3: 协议解析器 ⭐", "AI驱动3层架构/力科风格可视化"),
    ("07", "模块4: 代码生成器", "三级代码生成/驱动类型识别"),
    ("08", "模块5: 飞书集成", "对话式调用/零切换工作流"),
    ("09", "模块6: 知识库", "keyword-vault 10方向239关键词"),
    ("10", "技术架构与亮点", "Skill架构/MCP Server/降级策略"),
    ("11", "效果数据与价值", "效率提升/评审维度覆盖"),
]

for i, (num, title, desc) in enumerate(toc_items):
    row = i // 2
    col = i % 2
    x = 0.8 + col * 6.2
    y = 2.0 + row * 0.85

    # Number
    add_textbox(slide, x, y, 0.6, 0.5, num,
                font_size=20, bold=True, color=ACCENT_CYAN)
    # Title
    add_textbox(slide, x + 0.7, y, 5, 0.35, title,
                font_size=16, bold=True, color=WHITE)
    # Description
    add_textbox(slide, x + 0.7, y + 0.35, 5, 0.3, desc,
                font_size=11, color=DIM_GRAY)


# ============================================================
# SLIDE 3: BACKGROUND
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "课题背景", "USB/PD 驱动开发的现状与挑战")
add_page_number(slide, 3)

# Left panel - Industry context
add_card(slide, 0.8, 2.0, 5.5, 4.8)
add_textbox(slide, 1.2, 2.2, 5, 0.4, "行业背景",
            font_size=22, bold=True, color=ACCENT_CYAN)

bg_lines = [
    "• USB Type-C/PD 已成为消费电子充电接口的统一标准",
    "• PD 3.1 协议支持最高 240W 功率，复杂度大幅提升",
    "• 通信协议融合 (PD + USB4 + DP Alt Mode) 增加调试难度",
    "• 芯片厂商 (展锐/高通/MTK) 各自 TCPM 实现差异大",
    "",
    "• 底层驱动工程师花费大量时间在日志分析和 Bug 定位上",
    "• PD 协商涉及多消息交互时序，肉眼分析效率低",
    "• 故障案例分散，经验难以沉淀和复用",
]
add_multiline_textbox(slide, 1.2, 2.8, 5, 3.6, bg_lines,
                      font_size=14, color=LIGHT_GRAY, line_spacing=1.6)

# Right panel - Stats
add_card(slide, 7.0, 2.0, 5.5, 4.8)
add_textbox(slide, 7.4, 2.2, 5, 0.4, "目标用户",
            font_size=22, bold=True, color=ACCENT_GREEN)

stats = [
    "🎯 USB 底层驱动工程师",
    "🎯 PD 充电协议工程师",
    "🎯 硬件测试与验证工程师",
    "🎯 BSP/系统移植工程师",
    "",
    "📊 数据:",
    "• 日志分析日均耗时 2h+",
    "• PD 协商 Bug 平均定位周期 3-5天",
    "• 一次 USB 枚举失败需逐行比对 500+ 行 log",
    "• 跨团队经验传递依赖口头沟通",
]
add_multiline_textbox(slide, 7.4, 2.8, 5, 3.6, stats,
                      font_size=14, color=LIGHT_GRAY, line_spacing=1.6)


# ============================================================
# SLIDE 4: PAIN POINTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "痛点分析", "USB/PD 驱动开发的 4 大核心痛点")
add_page_number(slide, 4)

pain_points = [
    ("😰", "日志分散难聚合",
     "测试 log 散落在 Bugzilla / FTP / UNC 路径中\n需手动逐个下载比对，浪费大量时间"),
    ("🤯", "协议报文难理解",
     "PD 消息 Header 是 16-bit hex\nPDO 是 32-bit 编码，人工解码效率极低"),
    ("🐌", "Bug 定位周期长",
     "从收到 Bug 报告到根因定位\n平均需 3-5 天，涉及跨团队协调"),
    ("📚", "经验难以沉淀",
     "故障排查经验存于工程师脑中\n新人上手慢，重复踩坑"),
]

for i, (emoji, title, desc) in enumerate(pain_points):
    x = 0.6 + i * 3.15
    y = 2.0

    add_card(slide, x, y, 2.9, 4.5)

    # Emoji & Title
    add_textbox(slide, x + 0.3, y + 0.3, 2.3, 0.5, f"{emoji} {title}",
                font_size=18, bold=True, color=ACCENT_ORANGE)
    # Description
    add_textbox(slide, x + 0.3, y + 1.2, 2.3, 2.8, desc,
                font_size=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 5: SOLUTION OVERVIEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "解决方案: OWLMYCLAW 6 大模块", "AI 驱动的底层驱动全栈自动化研发助手")
add_page_number(slide, 5)

modules = [
    ("🐛", "Bug 智能助手", "查询/解析/搜索\nMCP Server\n离线模式",
     ACCENT_BLUE),
    ("🔍", "Log 分析器", "角色感知分析\n3类角色画像\n增强查看器",
     ACCENT_GREEN),
    ("📐", "协议解析器 ⭐", "AI驱动解析\n力科风格可视化\n任意格式兼容",
     ACCENT_CYAN),
    ("💻", "代码生成器", "三级代码生成\n驱动类型识别\n模板复用",
     ACCENT_ORANGE),
    ("💬", "飞书集成", "对话式调用\n卡片推送通知\n文件自动处理",
     RGBColor(0x9B, 0x59, 0xB6)),
    ("📚", "知识库", "10方向239关键词\nPD+USB协议知识\n故障案例库",
     ACCENT_RED),
]

for i, (icon, name, desc, color) in enumerate(modules):
    row = i // 3
    col = i % 3
    x = 0.6 + col * 4.15
    y = 2.0 + row * 2.65

    # Card with colored top border
    card = add_card(slide, x, y, 3.85, 2.35)

    # Colored accent bar on top of card
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y),
        Inches(3.85), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Icon & Name
    add_textbox(slide, x + 0.3, y + 0.3, 3.3, 0.45,
                f"{icon}  {name}", font_size=20, bold=True, color=color)
    # Description
    add_textbox(slide, x + 0.3, y + 0.9, 3.3, 1.2, desc,
                font_size=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 6: MODULE 1 - BUG ASSISTANT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "模块1: Bug 智能助手", "Bugzilla 查询/解析/搜索 · MCP Server · 离线模式")
add_page_number(slide, 6)

# Left column
add_card(slide, 0.6, 2.0, 5.8, 4.8)
add_textbox(slide, 1.0, 2.3, 5, 0.4, "核心能力",
            font_size=22, bold=True, color=ACCENT_BLUE)

features = [
    "🔍 Bug ID 查询 — 支持 5 大 Bugzilla 实例",
    "   (Kernel / Mozilla / GNOME / Unisoc / 本地)",
    "📊 智能摘要 — 自动提取关键信息生成摘要",
    "🔎 搜索功能 — 按关键词/严重度快速筛选",
    "📁 离线模式 — 预置模拟数据，演示/开发无网络可用",
    "📎 日志目录 — UNC+FTP 双通道自动列表",
    "📥 附件下载 — 智能识别 UNC/rayfile/SPCSS 渠道",
    "🔗 工作区 — 一键创建 + Bug Z 模板 + URL 快捷方式",
    "🤖 MCP Server — 5 个工具供 AI 直接调用",
]
add_multiline_textbox(slide, 1.0, 2.9, 5.2, 3.5, features,
                      font_size=14, color=LIGHT_GRAY, line_spacing=1.5)

# Right column - Architecture
add_card(slide, 7.0, 2.0, 5.5, 4.8)
add_textbox(slide, 7.4, 2.3, 5, 0.4, "Skill 拆分架构",
            font_size=22, bold=True, color=ACCENT_GREEN)

arch = [
    "原 bug-assistant (984行) 拆分为:",
    "",
    "📦 bug-core (320行)",
    "   Bug 查询 · 搜索 · 解析 · CLI",
    "",
    "📦 bug-logs (370行)",
    "   日志目录 · 结构化解析 · 附件下载",
    "",
    "📦 bug-workspace (230行)",
    "   工作区 · rayfile · 快捷方式",
    "",
    "🤖 MCP Server",
    "   统一入口 · 5工具 · 三级降级导入",
]
add_multiline_textbox(slide, 7.4, 2.9, 5, 3.5, arch,
                      font_size=13, color=LIGHT_GRAY, line_spacing=1.4)


# ============================================================
# SLIDE 7: MODULE 2 - LOG ANALYZER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "模块2: 角色感知日志分析器", "3 类角色画像 · 异常检测 · 时间线重建 · 增强查看器")
add_page_number(slide, 7)

# Role cards
roles = [
    ("⚡ PD 充电工程师", "关键词: PD, VBUS, SRC_CAP\nSNK_CAP, REQUEST, PDO\n关注: PD协商流程, 电压变化\nPDO匹配, PR_SWAP/DR_SWAP",
     ACCENT_BLUE),
    ("🔌 USB 驱动工程师", "关键词: Device Descriptor\nConfig Descriptor, Endpoint\n枚举, Reset, xhci\n关注: 枚举流程, 描述符读取",
     ACCENT_GREEN),
    ("🧪 测试工程师", "关键词: ERROR, WARN, FAIL\nTIMEOUT, Reset, Recovery\n关注: 跨模块日志关联\n异常聚类, 回归测试对比",
     ACCENT_ORANGE),
]

for i, (name, desc, color) in enumerate(roles):
    x = 0.6 + i * 4.15
    card = add_card(slide, x, 2.0, 3.85, 2.8)
    add_textbox(slide, x + 0.3, 2.2, 3.3, 0.4, name,
                font_size=18, bold=True, color=color)
    add_textbox(slide, x + 0.3, 2.8, 3.3, 1.8, desc,
                font_size=13, color=LIGHT_GRAY)

# Bottom - Viewer features
add_textbox(slide, 0.6, 5.2, 12, 0.4, "增强版日志查看器 (log_viewer_v2)",
            font_size=20, bold=True, color=ACCENT_CYAN)

viewer_features = [
    "🎯 左侧问题方向选择器 — 点击切换自动加载关键词 | 👤 徽章标识角色来源",
    "📋 关键词快捷复制 — 正则 OR / 空格 / 换行 / 通配 / grep -E 五种格式",
    "➕ 方向可视化 CRUD — 新建 / 编辑 / 删除问题方向，实时生效",
    "🔑 关键词可视化 CRUD — 添加 / 编辑 / 删除关键词",
    "📦 单文件分享 — 所有数据内联 HTML，一个文件即可分享",
]
add_multiline_textbox(slide, 0.6, 5.6, 12, 1.5, viewer_features,
                      font_size=13, color=LIGHT_GRAY, line_spacing=1.5)


# ============================================================
# SLIDE 8: MODULE 3 - PROTOCOL PARSER (CORE)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "模块3⭐: 协议智能解析与可视化", "AI 驱动 3 层架构 · 力科协议分析仪风格 · 兼容任意日志格式")
add_page_number(slide, 8)

# Architecture layers
layers = [
    ("Layer 1", "渐进式提取器", "4级提取策略\nLevel 0: 特定格式 (高置信度)\nLevel 1: 关键词匹配\nLevel 2: 超宽匹配\nLevel 3: 裸传尾部兜底",
     ACCENT_BLUE),
    ("Layer 2", "AI 语义理解", "自动解码Header hex\n解析PDO结构\n识别状态机流转\n检测异常/超时/死锁",
     ACCENT_CYAN),
    ("Layer 3", "可视化渲染", "力科风格仪表盘\nV/A/W 实时数据\nSVG 时序图\nPDO 能力表\n异常红色标注",
     ACCENT_GREEN),
]

for i, (layer, title, desc, color) in enumerate(layers):
    x = 0.6 + i * 4.15
    y = 2.0

    # Layer label
    add_textbox(slide, x, y, 3.85, 0.35, f"{layer}: {title}",
                font_size=16, bold=True, color=color)
    # Arrow (except last)
    if i < 2:
        add_textbox(slide, x + 3.6, y + 0.05, 0.5, 0.3, "→",
                    font_size=20, bold=True, color=ACCENT_CYAN)

    card = add_card(slide, x, y + 0.5, 3.85, 3.5)
    add_textbox(slide, x + 0.3, y + 0.7, 3.3, 2.8, desc,
                font_size=13, color=LIGHT_GRAY)

# Bottom info
add_textbox(slide, 0.6, 6.2, 12, 0.4,
            "核心创新: 不依赖固定正则 → AI + 知识库驱动解析，日志格式千变万化也能正确解析",
            font_size=14, bold=True, color=ACCENT_ORANGE)


# ============================================================
# SLIDE 9: MODULE 3 - SUPPORTED FORMATS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "协议解析器 — 支持的格式与能力", "Level 0 特定格式 + Level 1-3 兜底策略")
add_page_number(slide, 9)

# Level 0 formats
add_card(slide, 0.6, 2.0, 12.1, 4.5)
add_textbox(slide, 1.0, 2.2, 11, 0.4, "Level 0: 特定格式 (高置信度)",
            font_size=20, bold=True, color=ACCENT_GREEN)

formats_col1 = [
    "🟢 展锐 SC27XX TCPM/PD",
    "   (sprd_tcpm_log / sprd_pd_log)",
    "🟢 通用 PD Protocol [pd_protocol]",
    "🟢 通用 CC Protocol [cc_protocol]",
    "🟢 高通 PMIC PD [pmic_pd] / [qcom_pd]",
]
formats_col2 = [
    "🟢 通用 USB Core [usb_core]",
    "🟢 MTK TCPC [mtk_tcpc]",
    "🟢 FUSB302/TCPCI [fusb302] / [tcpci]",
    "🟡 Level 1-3: 关键词/超宽/裸传兜底",
]

add_multiline_textbox(slide, 1.0, 2.8, 5.5, 3.0, formats_col1,
                      font_size=14, color=LIGHT_GRAY, line_spacing=1.5)
add_multiline_textbox(slide, 6.8, 2.8, 5.5, 3.0, formats_col2,
                      font_size=14, color=LIGHT_GRAY, line_spacing=1.5)

# Test coverage
add_textbox(slide, 0.6, 5.0, 12, 0.4,
            "✅ 60 tests 全覆盖 · 支持远程测试控制 · 多次迭代至 v3.0",
            font_size=16, bold=True, color=ACCENT_GREEN)

# Bottom
add_card(slide, 0.6, 5.6, 12.1, 1.5)
add_textbox(slide, 1.0, 5.7, 11, 0.4, "架构演进历程",
            font_size=16, bold=True, color=ACCENT_CYAN)
evolution = "v1.0 正则匹配 → v2.0 CC解析器+增强可视化+60用例测试套件 → v3.0 AI驱动架构+渐进式提取器 (当前)"
add_textbox(slide, 1.0, 6.2, 11, 0.6, evolution, font_size=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 10: MODULE 4 - CODE GENERATOR
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "模块4: 代码辅助生成器", "三级分级代码生成 · 驱动类型智能识别 · 从分析到修复闭环")
add_page_number(slide, 10)

levels = [
    ("L1", "快速调试", "5分钟可用",
     "pr_debug()/pr_info() 日志\n简单条件判断 + 变量打印\n不改变程序逻辑",
     ACCENT_GREEN),
    ("L2", "完整诊断", "15分钟可用",
     "条件断点 + 状态收集 + 自动上报\n完整时序追踪\n关键状态机状态打印",
     ACCENT_BLUE),
    ("L3", "修复建议", "生产级",
     "补丁级代码 + 边界条件处理\n回归测试建议\n性能和安全影响分析",
     ACCENT_ORANGE),
]

for i, (level, name, tag, desc, color) in enumerate(levels):
    x = 0.6 + i * 4.15
    y = 2.0

    # Level badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 0.3), Inches(y),
        Inches(0.7), Inches(0.7)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = level
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + 1.1, y + 0.05, 2.5, 0.35,
                f"{name} ({tag})", font_size=18, bold=True, color=color)
    add_textbox(slide, x + 0.3, y + 1.0, 3.3, 1.5, desc,
                font_size=13, color=LIGHT_GRAY)

# Driver types
add_card(slide, 0.6, 4.2, 12.1, 2.5)
add_textbox(slide, 1.0, 4.4, 11, 0.35, "驱动类型智能识别",
            font_size=18, bold=True, color=ACCENT_CYAN)

driver_types = [
    ("PD/TCPM", "pd_protocol, VBUS\nPD_SRC_CAP"),
    ("USB Host/Device", "usb_core, Enumeration\nDescriptor"),
    ("I2C/PMIC", "i2c, charger, PMIC"),
    ("Thermal", "thermal, thermal_zone"),
    ("Kernel", "kernel, dmesg, Call Trace"),
]

for i, (dtype, kw) in enumerate(driver_types):
    x = 1.0 + i * 2.4
    add_textbox(slide, x, 4.95, 2.2, 0.3, dtype, font_size=14, bold=True,
                color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, 5.3, 2.2, 0.9, kw, font_size=11, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11: MODULE 5 - FEISHU INTEGRATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "模块5: 飞书 IM 集成", "对话式调用 · 卡片推送 · 文件自动处理 · 零切换工作流")
add_page_number(slide, 11)

# Left - Commands
add_card(slide, 0.6, 2.0, 5.8, 4.8)
add_textbox(slide, 1.0, 2.2, 5, 0.4, "对话指令",
            font_size=20, bold=True, color=ACCENT_CYAN)

cmds = [
    "/bug 查询    → 查询 Bug 信息",
    "/log 分析    → 上传日志 → 自动分析",
    "/code 生成   → 生成 debug 代码",
    "/protocol 解析 → 协议解析+可视化",
    "/bug 列表    → 查询 Bug 列表",
    "",
    "组合指令:",
    "/bug 分析 BUG-20261",
    "  → 查询Bug + 自动分析附件log",
    "/log 诊断 log.txt",
    "  → 日志分析 + 自动生成debug代码",
]
add_multiline_textbox(slide, 1.0, 2.8, 5.2, 3.5, cmds,
                      font_size=14, color=LIGHT_GRAY, line_spacing=1.5)

# Right - Card Messages
add_card(slide, 7.0, 2.0, 5.5, 4.8)
add_textbox(slide, 7.4, 2.2, 5, 0.4, "飞书卡片推送",
            font_size=20, bold=True, color=ACCENT_GREEN)

cards = [
    "📬 Bug 变更通知卡片",
    "   · 状态变更实时推送",
    "   · 一键查看详情/分析附件",
    "",
    "📊 日志分析结果卡片",
    "   · 异常统计摘要",
    "   · 严重程度标注",
    "   · 一键跳转生成debug代码",
    "",
    "📅 定时任务",
    "   · Bug 变更巡检 (每4小时)",
    "   · 每日 Bug 报告 (工作日9:00)",
]
add_multiline_textbox(slide, 7.4, 2.8, 5, 3.5, cards,
                      font_size=13, color=LIGHT_GRAY, line_spacing=1.4)


# ============================================================
# SLIDE 12: MODULE 6 - KNOWLEDGE BASE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "模块6: 知识库与关键词系统", "keyword-vault · 10方向239关键词 · PD/USB协议知识 · 故障案例库")
add_page_number(slide, 12)

# Keyword Vault
add_card(slide, 0.6, 2.0, 7.5, 4.8)
add_textbox(slide, 1.0, 2.2, 7, 0.4, "keyword-vault 关键词收纳系统",
            font_size=20, bold=True, color=ACCENT_CYAN)

directions_col1 = [
    "📋 协议超时问题 (21词)",
    "📋 功率/电压不匹配 (29词)",
    "📋 USB 枚举失败 (28词)",
    "📋 连接稳定性问题 (31词)",
    "📋 角色/方向切换问题 (19词)",
]
directions_col2 = [
    "📋 温度/热管理问题 (19词)",
    "📋 固件/启动问题 (27词)",
    "👤 PD工程师关注点 (19词)",
    "👤 USB工程师关注点 (24词)",
    "👤 测试工程师关注点 (22词)",
]

add_multiline_textbox(slide, 1.0, 2.8, 3.5, 3.5, directions_col1,
                      font_size=13, color=LIGHT_GRAY, line_spacing=1.5)
add_multiline_textbox(slide, 4.8, 2.8, 3.5, 3.5, directions_col2,
                      font_size=13, color=LIGHT_GRAY, line_spacing=1.5)

# Right - Knowledge base
add_card(slide, 8.5, 2.0, 4.2, 4.8)
add_textbox(slide, 8.8, 2.2, 3.5, 0.4, "知识库",
            font_size=20, bold=True, color=ACCENT_GREEN)

kb = [
    "📖 PD 3.1 协议规范",
    "   · 消息类型/PDO结构",
    "   · 状态机/时序要求",
    "",
    "📖 USB 枚举流程",
    "   · 描述符结构",
    "   · 传输类型/速率",
    "",
    "📖 故障案例库",
    "   · 5-10个典型故障案例",
    "   · 含log+分析+修复代码",
    "",
    "📖 提示词模板库",
    "   · 多场景优化提示词",
]
add_multiline_textbox(slide, 8.8, 2.8, 3.8, 3.5, kb,
                      font_size=12, color=LIGHT_GRAY, line_spacing=1.3)


# ============================================================
# SLIDE 13: TECHNICAL ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "技术架构与设计亮点", "Skill 架构 · MCP Server · 三级降级策略 · 模块解耦")
add_page_number(slide, 13)

# Architecture Diagram (text-based)
add_card(slide, 0.6, 2.0, 5.8, 2.5)
add_textbox(slide, 1.0, 2.2, 5, 0.4, "分层架构设计",
            font_size=20, bold=True, color=ACCENT_CYAN)
arch_text = [
    "📥 输入层:  文件上传 / Bug链接 / 飞书消息 / 手动粘贴",
    "🧠 分析层:  log-analyzer / protocol-parser / code-generator",
    "🗄️ 数据层:  keyword-vault / 知识库 / 故障案例 / 角色画像",
    "📤 输出层:  可视化HTML / 分析报告 / 飞书卡片 / MCP工具",
]
add_multiline_textbox(slide, 1.0, 2.8, 5.2, 1.5, arch_text,
                      font_size=14, color=LIGHT_GRAY, line_spacing=1.5)

# MCP Server
add_card(slide, 7.0, 2.0, 5.5, 2.5)
add_textbox(slide, 7.4, 2.2, 5, 0.4, "MCP Server (Model Context Protocol)",
            font_size=18, bold=True, color=ACCENT_GREEN)
mcp = [
    "🤖 5 个 MCP 工具供 AI 直接调用:",
    "   bug_info · bug_comments · bugs_search",
    "   bug_workspace · bug_logs",
    "📦 统一入口 · 模块自动发现",
    "🔗 三级降级导入策略",
]
add_multiline_textbox(slide, 7.4, 2.8, 5, 1.5, mcp,
                      font_size=13, color=LIGHT_GRAY, line_spacing=1.4)

# Design highlights
add_card(slide, 0.6, 4.8, 12.1, 2.2)
add_textbox(slide, 1.0, 5.0, 11, 0.35, "设计亮点",
            font_size=20, bold=True, color=ACCENT_CYAN)

highlights = [
    ("渐进式提取", "4级提取策略\n从特定格式到未知兜底", ACCENT_BLUE),
    ("AI语义理解", "不依赖正则\n知识库驱动正确解析", ACCENT_CYAN),
    ("三级降级", "try import\n→ sys.path\n→ 独立运行", ACCENT_GREEN),
    ("模块拆分", "984行→3独立Skill\n零耦合正交功能", ACCENT_ORANGE),
    ("单文件分享", "HTML内联数据\n无需服务器即可分享", RGBColor(0x9B, 0x59, 0xB6)),
]

for i, (name, desc, color) in enumerate(highlights):
    x = 1.0 + i * 2.3
    add_textbox(slide, x, 5.5, 2, 0.3, name, font_size=14, bold=True,
                color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, 5.85, 2, 1.0, desc, font_size=11, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 14: EFFECT DATA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "效果数据与量化指标", "效率提升 · 测试覆盖 · 模块完整度")
add_page_number(slide, 14)

metrics = [
    ("80%+", "日志分析时效提升", "从 2h+ 到 <30min\nAI自动提取关键信息", ACCENT_GREEN),
    ("60%+", "Bug定位加速", "从 3-5天 到 <2天\n自动关联log+协议分析", ACCENT_BLUE),
    ("239", "关键词沉淀", "10 个问题方向\n知识可复用可传承", ACCENT_CYAN),
    ("9", "独立 Skill", "完整功能覆盖\n模块化可插拔", ACCENT_ORANGE),
]

for i, (num, title, desc, color) in enumerate(metrics):
    x = 0.6 + i * 3.15
    y = 2.0
    add_card(slide, x, y, 2.9, 2.5)
    add_textbox(slide, x + 0.3, y + 0.2, 2.3, 0.7, num,
                font_size=36, bold=True, color=color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, y + 1.0, 2.3, 0.4, title,
                font_size=14, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, y + 1.4, 2.3, 0.8, desc,
                font_size=11, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

# Test coverage
add_card(slide, 0.6, 4.8, 5.8, 2.2)
add_textbox(slide, 1.0, 5.0, 5, 0.35, "测试覆盖",
            font_size=18, bold=True, color=ACCENT_GREEN)
test_info = [
    "📊 protocol-parser: 60 tests 100%通过",
    "📊 bug-core + bug-logs + bug-workspace: 136+ tests 100%通过",
    "📊 keyword-vault + log-analyzer: 完整测试套件",
    "📊 总计: 200+ 测试用例 · 10+ 测试组",
]
add_multiline_textbox(slide, 1.0, 5.5, 5.2, 1.3, test_info,
                      font_size=13, color=LIGHT_GRAY, line_spacing=1.4)

# Module completeness
add_card(slide, 7.0, 4.8, 5.5, 2.2)
add_textbox(slide, 7.4, 5.0, 5, 0.35, "模块完整度",
            font_size=18, bold=True, color=ACCENT_BLUE)
completeness = [
    "✅ bug-assistant (拆分为3)  — v2.0 下载+编排",
    "✅ log-analyzer              — v2.1 增强查看器",
    "✅ protocol-parser ⭐        — v3.0 AI驱动",
    "✅ keyword-vault             — v1.0 10方向239词",
    "✅ code-generator            — v1.0 三级生成",
    "✅ feishu-integration        — v1.0 飞书对话",
]
add_multiline_textbox(slide, 7.4, 5.5, 5, 1.3, completeness,
                      font_size=13, color=LIGHT_GRAY, line_spacing=1.3)


# ============================================================
# SLIDE 15: VALUE PROPOSITION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "课题价值", "可应用性 · 可推广性 · 可提效性")
add_page_number(slide, 15)

value_cards = [
    ("🎯 可应用性", ACCENT_BLUE, [
        "✅ 覆盖完整 USB/PD 驱动开发调试闭环",
        "✅ 6 大模块可直接用于日常工作",
        "✅ 支持真实 Bugzilla + 离线演示双模式",
        "✅ 兼容展锐/高通/MTK 多平台日志格式",
        "✅ 从分析到代码生成一站式工作流",
    ]),
    ("🚀 可推广性", ACCENT_GREEN, [
        "✅ 插件化 Skill 架构，模块独立可插拔",
        "✅ keyword-vault 知识底座可适配其他领域",
        "✅ MCP Server 标准化接口，可接入其他 AI 平台",
        "✅ 飞书集成模式可直接复制到其他团队",
        "✅ 三级降级策略保证独立运行不依赖外部",
    ]),
    ("⚡ 可提效性", ACCENT_ORANGE, [
        "✅ 日志分析: 2h → 30min (提效 75%)",
        "✅ Bug 定位: 3-5天 → 1-2天 (提效 60%)",
        "✅ 关键词系统: 经验从人脑 → 可复用知识库",
        "✅ 代码生成: 手动编写 → 一键生成三级代码",
        "✅ 飞书集成: 多工具切换 → 一个对话窗口",
    ]),
]

for i, (title, color, items) in enumerate(value_cards):
    x = 0.4 + i * 4.2
    card = add_card(slide, x, 2.0, 4.0, 5.0)
    # Colored top bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(2.0),
        Inches(4.0), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide, x + 0.3, 2.3, 3.4, 0.5, title,
                font_size=22, bold=True, color=color)
    add_multiline_textbox(slide, x + 0.3, 2.9, 3.4, 3.8, items,
                          font_size=13, color=LIGHT_GRAY, line_spacing=1.7)


# ============================================================
# SLIDE 16: SUMMARY & OUTLOOK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title_bar(slide, "总结与展望", "Summary & Future Work")
add_page_number(slide, 16)

# Summary
add_card(slide, 0.6, 2.0, 5.8, 4.8)
add_textbox(slide, 1.0, 2.2, 5, 0.4, "项目总结",
            font_size=22, bold=True, color=ACCENT_CYAN)

summary = [
    "OWLMYCLAW 是一个面向 USB/PD 底层驱动",
    "开发工程师的 AI 辅助工具集，通过 9 个独立",
    "Skill 覆盖从 Bug 接收到代码修复的完整链路。",
    "",
    "🔑 核心创新点:",
    "• AI + 知识库驱动协议解析",
    "  不依赖固定正则，兼容任意日志格式",
    "• 渐进式提取器 4 级降级策略",
    "  从特定格式到完全未知都能处理",
    "• keyword-vault 知识底座",
    "  让调试经验可沉淀、可复用、可传承",
    "• MCP Server 标准化接口",
    "  模块解耦，易于扩展和集成",
]
add_multiline_textbox(slide, 1.0, 2.8, 5.2, 3.8, summary,
                      font_size=14, color=LIGHT_GRAY, line_spacing=1.4)

# Future work
add_card(slide, 7.0, 2.0, 5.5, 4.8)
add_textbox(slide, 7.4, 2.2, 5, 0.4, "未来展望",
            font_size=22, bold=True, color=ACCENT_GREEN)

future = [
    "📋 短期 (1-3月)",
    "• 扩展更多芯片平台日志格式支持",
    "• keyword-vault 自动学习能力增强",
    "• 故障案例库充实至 50+",
    "",
    "📋 中期 (3-6月)",
    "• 集成 USB4 / Thunderbolt 协议解析",
    "• 多维度回归测试对比分析",
    "• 企业微信/钉钉渠道扩展",
    "",
    "📋 长期 (6-12月)",
    "• 扩展到其他驱动领域 (I2C/SPI/PCIe)",
    "• 自动化测试流程集成",
    "• 团队级知识库共建共享",
]
add_multiline_textbox(slide, 7.4, 2.8, 5, 3.8, future,
                      font_size=12, color=LIGHT_GRAY, line_spacing=1.3)


# ============================================================
# SLIDE 17: ENDING / THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# Decorative bar
bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    SLIDE_WIDTH, Inches(0.15)
)
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# Thank you
add_textbox(slide, 1.0, 1.8, 11, 1.0, "感谢观看",
            font_size=60, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 2.8, 11, 0.6, "Thank You",
            font_size=28, color=ACCENT_CYAN,
            alignment=PP_ALIGN.CENTER)

# Divider
div = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(5.5), Inches(3.6),
    Inches(2.3), Inches(0.04)
)
div.fill.solid()
div.fill.fore_color.rgb = ACCENT_CYAN
div.line.fill.background()

# Info
add_textbox(slide, 1.0, 4.0, 11, 0.5, "OWLMYCLAW — 底层驱动全栈自动化研发助手",
            font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_multiline_textbox(slide, 1.0, 4.8, 11, 1.5, [
    "2026 火山杯 · 紫光展锐 AI 创新大赛",
    "研发与创新赛道",
    "",
    "github.com/owlm00n/bugzilla_assistent",
], font_size=16, color=DIM_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

# Bottom bar
bar2 = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(7.35),
    SLIDE_WIDTH, Inches(0.15)
)
bar2.fill.solid()
bar2.fill.fore_color.rgb = ACCENT_BLUE
bar2.line.fill.background()


# ============================================================
# SAVE
# ============================================================
os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
prs.save(OUTPUT_PATH)
print(f"✅ PPT generated successfully: {OUTPUT_PATH}")
print(f"   Slides: {len(prs.slides)}")
